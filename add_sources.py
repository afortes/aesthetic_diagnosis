import os
import time
import hashlib
import json
from dotenv import load_dotenv
from pinecone import Pinecone
from pathlib import Path

load_dotenv()

# --- CONFIGURACIÓN ---
DOCS_DIR = "./fuentes"
NAMESPACE = "example-namespace"
INDEX_NAME = "aesthetic-knowledge"
CHUNK_SIZE = 1000       # tamaño máximo de cada chunk (caracteres) para texto/pptx
CHUNK_OVERLAP = 200     # solapamiento entre chunks
BATCH_SIZE = 64         # tamaño del lote para subir a Pinecone
TOKEN_LIMIT_PER_MINUTE = 240000 # Límite de seguridad
WAIT_TIME_SECONDS = 65  # 1 minuto y 5 segundos

pc = Pinecone(api_key=os.environ["PINECONE_API_KEY"])

if not pc.has_index(INDEX_NAME):
    print(f"🔧 Índice '{INDEX_NAME}' no encontrado. Creándolo...")
    pc.create_index_for_model(
        name=INDEX_NAME,
        cloud="aws",
        region="us-east-1",
        embed={
            "model": "llama-text-embed-v2",
            "field_map": {"text": "chunk_text"}
        }
    )
    print(f"✅ Índice '{INDEX_NAME}' creado correctamente.")

index = pc.Index(host=os.environ["PINECONE_INDEX_HOST"])

# --- FUNCIONES AUXILIARES DE LECTURA ---

def read_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_md(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_xlsx(path):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        all_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = [row for row in ws.iter_rows(values_only=True) if any(c is not None for c in row)]
            if not rows:
                continue

            # Buscar la fila de cabeceras: la primera fila donde todos los valores no-None son strings
            header_idx = 0
            for idx, row in enumerate(rows):
                non_empty = [c for c in row if c is not None]
                if non_empty and all(isinstance(c, str) for c in non_empty):
                    header_idx = idx
                    break

            headers = [str(c).strip() if c is not None else "" for c in rows[header_idx]]
            data_rows = rows[header_idx + 1:]

            all_text.append(f"[Hoja: {sheet_name}]")
            for row in data_rows:
                parts = []
                for header, value in zip(headers, row):
                    if value is not None and str(value).strip():
                        parts.append(f"{header}: {value}")
                if parts:
                    all_text.append(" | ".join(parts))

        wb.close()
        return "\n".join(all_text)
    except Exception as e:
        print(f"⚠️ Error al leer XLSX {path}: {e}")
        return ""

def read_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"⚠️ Error al leer DOCX {path}: {e}")
        return ""

def read_pdf(path):
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        pages = [p.extract_text() or "" for p in reader.pages]
        return "\n".join(pages)
    except Exception as e:
        print(f"⚠️ Error al leer PDF {path}: {e}")
        return ""

def read_pptx(path):
    """Extrae texto de un archivo PowerPoint estructurado por diapositivas."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slide_chunks = []
        
        for i, slide in enumerate(prs.slides):
            slide_text_elements = []
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
                
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text != title:
                        slide_text_elements.append(text)
                        
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text_elements.append("\n[Notas del Orador]:\n" + notes)
                    
            full_text = "\n".join(slide_text_elements)
            if not full_text and not title:
                continue
                
            chunk_content = f"--- Diapositiva {i+1} ---\n"
            if title:
                chunk_content += f"Título: {title}\n"
            chunk_content += f"Contenido:\n{full_text}"
            
            slide_chunks.append({
                "text": chunk_content,
                "metadata": {"slide_number": i + 1, "slide_title": title[:50] if title else "Sin título"}
            })
        return slide_chunks
    except Exception as e:
        print(f"⚠️ Error al leer PPTX {path}: {e}")
        return []

def read_mp4(path):
    """
    Transcribe con Whisper (usando caché .json) y devuelve chunks con overlap.
    Whisper lee el MP4 directamente via ffmpeg, sin extracción previa de audio.
    """
    try:
        import whisper

        cache_file = path.with_suffix('.json')

        # 1. Comprobar caché
        if cache_file.exists():
            print(f"📦 Usando transcripción en caché para {path.name}...")
            with open(cache_file, "r", encoding="utf-8") as f:
                segments = json.load(f)
        else:
            print(f"🎙️ Transcribiendo {path.name} (esto puede tardar)...")
            model = whisper.load_model("small")
            result = model.transcribe(str(path), language="es")
            segments = result["segments"]

            with open(cache_file, "w", encoding="utf-8") as f:
                json.dump(segments, f, ensure_ascii=False, indent=4)
            print(f"💾 Transcripción guardada en caché: {cache_file.name}")

            # Guardar versión legible para revisión humana
            txt_file = path.with_suffix('.txt')
            with open(txt_file, "w", encoding="utf-8") as f:
                f.write(f"TRANSCRIPCIÓN: {path.name}\n")
                f.write("=" * 60 + "\n\n")
                for seg in segments:
                    start = int(seg["start"])
                    end = int(seg["end"])
                    timestamp = f"[{start // 60:02d}:{start % 60:02d} - {end // 60:02d}:{end % 60:02d}]"
                    f.write(f"{timestamp} {seg['text'].strip()}\n")
            print(f"📄 Transcripción legible guardada en: {txt_file.name}")

        # 2. Chunking con overlap basado en tiempo
        # Construimos primero la lista completa de (texto, start, end) por segmento
        audio_chunks = []
        i = 0
        while i < len(segments):
            chunk_text = ""
            chunk_start = segments[i]["start"]
            chunk_end = segments[i]["end"]
            j = i

            # Acumular segmentos hasta alcanzar CHUNK_SIZE
            while j < len(segments) and len(chunk_text) < CHUNK_SIZE:
                chunk_text += segments[j]["text"] + " "
                chunk_end = segments[j]["end"]
                j += 1

            audio_chunks.append({
                "text": f"--- Transcripción de Audio [{chunk_start:.1f}s - {chunk_end:.1f}s] ---\n{chunk_text.strip()}",
                "metadata": {
                    "start_time": round(chunk_start, 2),
                    "end_time": round(chunk_end, 2)
                }
            })

            # Retroceder CHUNK_OVERLAP caracteres equivalentes en segmentos para el overlap
            overlap_text = ""
            k = j - 1
            while k >= i and len(overlap_text) < CHUNK_OVERLAP:
                overlap_text = segments[k]["text"] + " " + overlap_text
                k -= 1
            # El siguiente chunk empieza desde el segmento donde comienza el overlap
            i = max(k + 1, i + 1)

        return audio_chunks

    except Exception as e:
        print(f"⚠️ Error al procesar MP4 {path}: {e}")
        return []

def chunk_text(text, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Aplica particionado clásico por caracteres (usado para txt, pdf y pptx largos)."""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap if end < len(text) else end
    return chunks

def read_file(path: Path):
    ext = path.suffix.lower()
    if ext == ".txt":
        return read_txt(path)
    elif ext == ".md":
        return read_md(path)
    elif ext == ".pdf":
        return read_pdf(path)
    elif ext == ".docx":
        return read_docx(path)
    elif ext == ".xlsx":
        return read_xlsx(path)
    elif ext == ".pptx":
        return read_pptx(path)
    elif ext == ".mp4":
        return read_mp4(path)
    else:
        return ""

# --- REGISTRO DE ARCHIVOS YA INGESTADOS ---

INGESTED_LOG = Path(DOCS_DIR) / ".ingested.json"

def load_ingested_log():
    if INGESTED_LOG.exists():
        with open(INGESTED_LOG, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def mark_as_ingested(log, file):
    log[file.name] = {
        "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        "size_bytes": file.stat().st_size
    }
    with open(INGESTED_LOG, "w", encoding="utf-8") as f:
        json.dump(log, f, ensure_ascii=False, indent=4)

# --- PROCESO PRINCIPAL ---

def main():
    all_records = []
    ingested_log = load_ingested_log()
    files_to_process = []

    for file in Path(DOCS_DIR).glob("*"):
        if file.name.startswith(".") or file.name.startswith("~"):
            continue
        ext = file.suffix.lower()
        if ext not in [".txt", ".pdf", ".md", ".docx", ".xlsx", ".pptx", ".mp4"]:
            continue

        # Saltar .txt que son sidecars de transcripción de audio
        if ext == ".txt" and file.with_suffix(".mp4").exists():
            continue

        if file.name in ingested_log:
            print(f"⏭️  {file.name} ya ingestado ({ingested_log[file.name]['ingested_at']}), saltando.")
            continue

        files_to_process.append(file)

    if not files_to_process:
        print("✅ Todos los archivos ya están ingestados. No hay nada nuevo que procesar.")
        return

    for file in files_to_process:
        ext = file.suffix.lower()
        print(f"\nProcesando {file.name}...")
        content = read_file(file)
        file_records = []

        # 1. Procesamiento para archivos de texto plano (PDF, TXT, MD, DOCX)
        if isinstance(content, str):
            if not content.strip():
                print(f"⚠️ {file.name} vacío o sin texto legible.")
                continue

            chunks = chunk_text(content)
            for i, chunk in enumerate(chunks):
                file_records.append({
                    "_id": hashlib.md5(f"{file.stem}-chunk-{i}".encode()).hexdigest(),
                    "chunk_text": chunk,
                    "category": "general",
                    "source_file": file.name
                })

        # 2. Procesamiento para archivos estructurados (PPTX y MP4 devuelven listas)
        elif isinstance(content, list):
            if not content:
                print(f"⚠️ {file.name} vacío o sin contenido procesable.")
                continue

            if ext == ".pptx":
                for slide_data in content:
                    slide_text = slide_data["text"]
                    slide_meta = slide_data["metadata"]
                    sub_chunks = chunk_text(slide_text) if len(slide_text) > CHUNK_SIZE else [slide_text]
                    for i, chunk in enumerate(sub_chunks):
                        file_records.append({
                            "_id": hashlib.md5(f"{file.stem}-slide{slide_meta['slide_number']}-{i}".encode()).hexdigest(),
                            "chunk_text": chunk,
                            "category": "presentacion",
                            "source_file": file.name,
                            "slide_number": slide_meta["slide_number"],
                            "slide_title": slide_meta["slide_title"]
                        })

            elif ext == ".mp4":
                for i, chunk_data in enumerate(content):
                    file_records.append({
                        "_id": hashlib.md5(f"{file.stem}-time{chunk_data['metadata']['start_time']}-{i}".encode()).hexdigest(),
                        "chunk_text": chunk_data["text"],
                        "category": "video",
                        "source_file": file.name,
                        "start_time": chunk_data["metadata"]["start_time"],
                        "end_time": chunk_data["metadata"]["end_time"]
                    })

        all_records.extend(file_records)

    # --- SUBIDA A PINECONE ---
    if not all_records:
        print("⚠️ No se encontraron archivos válidos o procesables.")
        return

    print(f"\nSubiendo {len(all_records)} chunks a Pinecone en lotes de {BATCH_SIZE}...")
    tokens_used_in_minute = 0

    # Agrupamos los records por archivo para marcar cada uno al terminar
    records_by_file = {}
    for record in all_records:
        records_by_file.setdefault(record["source_file"], []).append(record)

    for filename, records in records_by_file.items():
        file_path = Path(DOCS_DIR) / filename
        for i in range(0, len(records), BATCH_SIZE):
            batch = records[i: i + BATCH_SIZE]
            batch_text_len = sum(len(r["chunk_text"]) for r in batch)
            estimated_tokens = int(batch_text_len / 3.5)

            if tokens_used_in_minute + estimated_tokens >= TOKEN_LIMIT_PER_MINUTE:
                print(f"⚠️ Límite de tokens/min alcanzado. Esperando {WAIT_TIME_SECONDS}s...")
                time.sleep(WAIT_TIME_SECONDS)
                tokens_used_in_minute = 0
                print("⏳ Reanudando subida...")

            print(f"  Subiendo lote {i // BATCH_SIZE + 1} de {filename} ({len(batch)} registros)...")
            index.upsert_records(NAMESPACE, batch)
            tokens_used_in_minute += estimated_tokens

        mark_as_ingested(ingested_log, file_path)
        print(f"✅ {filename} completado y registrado.")

    print("\n✅ Carga completada correctamente.")

if __name__ == "__main__":
    main()